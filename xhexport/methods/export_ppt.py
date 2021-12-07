import json
import time
from os import path
from selenium import webdriver
from selenium.webdriver.common.keys import Keys
from selenium.webdriver.chrome.options import Options
from pptx import Presentation
from pptx.util import Inches
from xhexport import config
from xhexport.utils import logger

log = logger('export-ppt')

chrome_options = Options()
chrome_options.add_argument('--headless')
chrome_options.add_argument('--log-level=2')
chrome_options.add_argument('--window-size=4096x3072')


def hide_float(driver):
    return driver.execute_script(f'''
        document.getElementsByClassName('page')[0].style.display = 'none'
        document.getElementsByClassName('slidePic')[0].style.display = 'none'
    ''')


def simulate_click(driver, times):
    return driver.execute_script(f'''
        for (let _ = 0; _ < {times}; _++) {{
            if (console.log("[xhexport]simu-click"), K.tl && K.tl["sp-1"] && -1 != K.tl["sp-1"].bj) K.tl["sp-1"].bu();
            K.tl && K.tl["sp-1"] && K.tl["sp-1"].ak != K.tl["sp-1"].au.length ? K.tl["sp-1"].dG(K.tl["sp-1"]) : bE(aJ.ak + 1)
        }}
        return _XH.actionList[_XH.actionList.length - 1]
    ''')


def goto_page(driver, page):
    return driver.execute_script(f'''
        bE({page - 1})
    ''')


def export_per_page(html_path, dist_path):
    driver = webdriver.Chrome(executable_path=config.chrome_driver,
                              chrome_options=chrome_options)
    driver.get("file://" + html_path)

    time.sleep(3)

    hide_float(driver)
    last_status = None
    while True:
        current_status = simulate_click(driver, 20)[:-1]
        if last_status == current_status:
            break
        else:
            last_status = current_status

    total_page = last_status[0]
    ppt = Presentation()

    for page in range(1, total_page + 1):
        img_path = path.abspath(path.join(path.dirname(dist_path), \
                                          './screen-%03d.png' % page))
        log(f'保存第 {page}/{total_page} 页到', img_path)
        goto_page(driver, page)
        driver.save_screenshot(img_path)
        slide = ppt.slides.add_slide(ppt.slide_layouts[0])
        slide.shapes.add_picture(img_path, 0, 0, width=Inches(10))

    ppt.save(dist_path)
    driver.close()
